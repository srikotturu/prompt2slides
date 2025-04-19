# PowerPoint Utility Functions using python-pptx
# Requires: pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE



import os


def create_presentation():
    """
    Creates a new PowerPoint presentation
    
    Returns:
        Presentation: A new PowerPoint presentation object
    """
    prs = Presentation()
    return prs


def add_slide(prs, layout_index=1, title=None, background_color=None):
    """
    Adds a new slide to the presentation
    
    Args:
        prs (Presentation): The presentation object
        layout_index (int): Index of slide layout to use (default: 1, Title and Content)
        title (str, optional): Slide title
        background_color (tuple, optional): RGB tuple for background color (e.g., (255, 255, 255))
    
    Returns:
        slide: The newly created slide
    """
    # Available layout indices (may vary by template):
    # 0 - Title Slide
    # 1 - Title and Content
    # 2 - Section Header
    # 3 - Two Content
    # 4 - Comparison
    # 5 - Title Only
    # 6 - Blank
    # 7 - Content with Caption
    # 8 - Picture with Caption
    
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title if provided
    if title and slide.shapes.title:
        slide.shapes.title.text = title
    
    # Set background color if provided
    if background_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*background_color)
    
    return slide


def add_text(slide, text, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(1), 
             font_size=Pt(18), font_name="Calibri", color=(0, 0, 0), bold=False, 
             italic=False, underline=False, alignment=PP_ALIGN.LEFT):
    """
    Adds text to a slide
    
    Args:
        slide: The slide object
        text (str): Text content
        left (Inches): Left position
        top (Inches): Top position
        width (Inches): Width of text box
        height (Inches): Height of text box
        font_size (Pt): Font size
        font_name (str): Font name
        color (tuple): RGB color tuple (e.g., (255, 0, 0) for red)
        bold (bool): Whether text should be bold
        italic (bool): Whether text should be italic
        underline (bool): Whether text should be underlined
        alignment: Text alignment (PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT)
    
    Returns:
        shape: The text box shape
    """
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment
    
    run = paragraph.runs[0]
    run.font.size = font_size
    run.font.name = font_name
    run.font.color.rgb = RGBColor(*color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    
    return text_box


def add_paragraph(slide, text, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(3),
                  font_size=Pt(16), font_name="Calibri", color=(0, 0, 0),
                  alignment=PP_ALIGN.LEFT, line_spacing=1.0):
    """
    Adds a paragraph (multi-line text) to a slide
    
    Args:
        slide: The slide object
        text (str): Paragraph text content
        left (Inches): Left position
        top (Inches): Top position
        width (Inches): Width of text box
        height (Inches): Height of text box
        font_size (Pt): Font size
        font_name (str): Font name
        color (tuple): RGB color tuple
        alignment: Text alignment (PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT)
        line_spacing (float): Line spacing multiplier
    
    Returns:
        shape: The text box shape
    """
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment
    paragraph.line_spacing = line_spacing
    
    run = paragraph.runs[0]
    run.font.size = font_size
    run.font.name = font_name
    run.font.color.rgb = RGBColor(*color)
    
    return text_box


def add_bullet_list(slide, items, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(3),
                    font_size=Pt(16), font_name="Calibri", color=(0, 0, 0),
                    level=0, bullet_character=None):
    """
    Adds a bulleted list to a slide
    
    Args:
        slide: The slide object
        items (list): List of text items
        left (Inches): Left position
        top (Inches): Top position
        width (Inches): Width of text box
        height (Inches): Height of text box
        font_size (Pt): Font size
        font_name (str): Font name
        color (tuple): RGB color tuple
        level (int): Indentation level (0 for top level)
        bullet_character (str, optional): Custom bullet character
    
    Returns:
        shape: The text box shape
    """
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Clear default paragraph if it exists
    if text_frame.paragraphs:
        text_frame.paragraphs[0].text = ""
    
    for i, item in enumerate(items):
        # Add a paragraph for each bullet item
        if i == 0 and text_frame.paragraphs:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = level
        
        # Set bullet properties
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = RGBColor(*color)
        
        # Custom bullet character if provided
        if bullet_character:
            p.bullet.character = bullet_character
        else:
            p.bullet.type = None  # Let PowerPoint handle default bullet
    
    return text_box


def add_image(slide, image_path, left=Inches(1), top=Inches(2), width=None, height=None):
    """
    Adds an image to a slide
    
    Args:
        slide: The slide object
        image_path (str): Path to the image file
        left (Inches): Left position
        top (Inches): Top position
        width (Inches, optional): Width of image (maintains aspect ratio if only width is specified)
        height (Inches, optional): Height of image (maintains aspect ratio if only height is specified)
    
    Returns:
        shape: The image shape
    """
    if width is None and height is None:
        # Default size if neither width nor height is specified
        width, height = Inches(4), Inches(3)
    
    image = slide.shapes.add_picture(image_path, left, top, width, height)
    return image


def add_chart(slide, chart_type, categories, data_series, left=Inches(1), top=Inches(2),
              width=Inches(8), height=Inches(4), chart_title=None):
    """
    Adds a chart to a slide
    
    Args:
        slide: The slide object
        chart_type (XL_CHART_TYPE): Type of chart (e.g., XL_CHART_TYPE.BAR_CLUSTERED)
        categories (list): List of category labels
        data_series (list): List of tuples (series_name, values)
        left (Inches): Left position
        top (Inches): Top position
        width (Inches): Width of chart
        height (Inches): Height of chart
        chart_title (str, optional): Chart title
    
    Returns:
        shape: The chart shape
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for series_name, values in data_series:
        chart_data.add_series(series_name, values)
    
    chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart
    
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
    
    return chart


def add_table(slide, data, left=Inches(1), top=Inches(2), width=None, height=None,
              column_widths=None, first_row_is_header=True):
    """
    Adds a table to a slide
    
    Args:
        slide: The slide object
        data (list): 2D array/list of data
        left (Inches): Left position
        top (Inches): Top position
        width (Inches, optional): Total width of table
        height (Inches, optional): Total height of table
        column_widths (list, optional): List of column widths (Inches)
        first_row_is_header (bool): Whether to format first row as header
    
    Returns:
        shape: The table shape
    """
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if rows == 0 or cols == 0:
        return None
    
    # Calculate dimensions if not provided
    if width is None:
        if column_widths:
            width = sum(width for width in column_widths)
        else:
            width = Inches(8)  # Default width
    
    if height is None:
        height = Inches(0.5 * rows)  # Rough estimate based on rows
    
    # Create table
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths if provided
    if column_widths and len(column_widths) == cols:
        for i, width in enumerate(column_widths):
            table.columns[i].width = width
    
    # Populate table with data
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            
            # Handle both simple strings and complex data with formatting
            if isinstance(cell_data, dict):
                text = cell_data.get('text', '')
                options = cell_data.get('options', {})
                
                # Apply text
                p = cell.text_frame.paragraphs[0]
                p.text = text
                
                # Apply formatting if specified
                if options:
                    run = p.runs[0]
                    if 'bold' in options:
                        run.font.bold = options['bold']
                    if 'italic' in options:
                        run.font.italic = options['italic']
                    if 'font_size' in options:
                        run.font.size = options['font_size']
                    if 'color' in options:
                        run.font.color.rgb = RGBColor(*options['color'])
            else:
                cell.text = str(cell_data)
    
    # Format header row if specified
    if first_row_is_header and rows > 0:
        for cell in table.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
    
    return table


def add_shape(slide, shape_type, left=Inches(2), top=Inches(2), width=Inches(2), height=Inches(2),
             fill_color=None, line_color=None, line_width=Pt(1)):
    """
    Adds a shape to a slide
    
    Args:
        slide: The slide object
        shape_type (MSO_SHAPE): Type of shape
        left (Inches): Left position
        top (Inches): Top position
        width (Inches): Width of shape
        height (Inches): Height of shape
        fill_color (tuple, optional): RGB fill color
        line_color (tuple, optional): RGB line color
        line_width (Pt): Line width
    
    Returns:
        shape: The shape object
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    if fill_color:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*fill_color)
    
    if line_color:
        line = shape.line
        line.color.rgb = RGBColor(*line_color)
        line.width = line_width
    
    return shape


def add_header_footer(prs, header_text=None, footer_text=None, slide_number=True, date=True):
    """
    Adds headers and footers to all slides
    
    Args:
        prs (Presentation): The presentation object
        header_text (str, optional): Header text
        footer_text (str, optional): Footer text
        slide_number (bool): Whether to show slide numbers
        date (bool): Whether to show date
    """
    # Note: Headers are not directly supported in python-pptx
    # We can only add footers, slide numbers, and dates
    
    for i, slide in enumerate(prs.slides):
        if hasattr(slide, 'has_header_footer') and slide.has_header_footer:
            if footer_text:
                slide.footer.text = footer_text
                slide.footer.visible = True
            
            if slide_number:
                slide.slide_number.visible = True
            
            if date:
                slide.date.visible = True
            
            # For header, we need to add a text box at the top of each slide
            if header_text and i > 0:  # Skip title slide
                add_text(slide, header_text, 
                        left=Inches(0.5), 
                        top=Inches(0.1), 
                        width=Inches(9), 
                        height=Inches(0.5),
                        font_size=Pt(10),
                        color=(100, 100, 100),
                        alignment=PP_ALIGN.RIGHT)


def save_presentation(prs, filename="presentation"):
    """
    Saves the presentation to a file
    
    Args:
        prs (Presentation): The presentation object
        filename (str): Name of the file (without extension)
    """
    if not filename.endswith('.pptx'):
        filename += '.pptx'
    
    prs.save(filename)


def create_complete_presentation(title, slides_content, options=None):
    """
    Creates a complete PowerPoint presentation with multiple slides
    
    Args:
        title (str): Presentation title
        slides_content (list): List of slide content dictionaries
        options (dict, optional): Presentation options
    
    Returns:
        Presentation: The completed presentation object
    """
    if options is None:
        options = {}
    
    prs = create_presentation()
    
    # Add title slide
    title_slide_color = options.get('title_slide_color', (0, 136, 204))  # Default blue
    title_slide = add_slide(prs, 0, background_color=title_slide_color)  # 0 is Title Slide layout
    
    add_text(title_slide, title,
             left=Inches(1), 
             top=Inches(2.5), 
             width=Inches(8), 
             height=Inches(1.5),
             font_size=Pt(44),
             color=(255, 255, 255),
             bold=True,
             alignment=PP_ALIGN.CENTER)
    
    if 'subtitle' in options:
        add_text(title_slide, options['subtitle'],
                left=Inches(1), 
                top=Inches(4), 
                width=Inches(8), 
                height=Inches(1),
                font_size=Pt(24),
                color=(255, 255, 255),
                alignment=PP_ALIGN.CENTER)
    
    # Add content slides
    for slide_content in slides_content:
        # Default to Title and Content layout (index 1)
        layout_index = slide_content.get('layout_index', 1)
        slide = add_slide(prs, layout_index, slide_content.get('title'))
        
        # Current vertical position tracker for sequential content placement
        current_top = Inches(1.5)
        
        # Add text content
        if 'text' in slide_content:
            text_height = Inches(len(slide_content['text']) / 100 + 0.5)  # Rough estimate
            add_paragraph(slide, slide_content['text'],
                         top=current_top,
                         height=text_height)
            current_top += text_height + Inches(0.2)  # Add spacing
        
        # Add bullet points
        if 'bullets' in slide_content:
            bullets_height = Inches(len(slide_content['bullets']) * 0.3 + 0.2)
            add_bullet_list(slide, slide_content['bullets'],
                           top=current_top,
                           height=bullets_height)
            current_top += bullets_height + Inches(0.2)
        
        # Add image
        if 'image' in slide_content:
            img_options = slide_content.get('image_options', {})
            img_top = img_options.get('top', current_top)
            img_left = img_options.get('left', Inches(1))
            img_width = img_options.get('width', Inches(4))
            img_height = img_options.get('height', None)
            
            add_image(slide, slide_content['image'], 
                     left=img_left, 
                     top=img_top,
                     width=img_width,
                     height=img_height)
            
            if img_height:
                current_top = img_top + img_height + Inches(0.2)
            else:
                # Estimate based on width if height not specified
                current_top += Inches(3) + Inches(0.2)
        
        # Add chart
        if 'chart_type' in slide_content and 'chart_categories' in slide_content and 'chart_data' in slide_content:
            chart_options = slide_content.get('chart_options', {})
            chart_top = chart_options.get('top', current_top)
            chart_left = chart_options.get('left', Inches(1))
            chart_width = chart_options.get('width', Inches(8))
            chart_height = chart_options.get('height', Inches(4))
            chart_title = chart_options.get('title')
            
            add_chart(slide, 
                     slide_content['chart_type'],
                     slide_content['chart_categories'],
                     slide_content['chart_data'],
                     left=chart_left,
                     top=chart_top,
                     width=chart_width,
                     height=chart_height,
                     chart_title=chart_title)
            
            current_top = chart_top + chart_height + Inches(0.2)
        
        # Add table
        if 'table_data' in slide_content:
            table_options = slide_content.get('table_options', {})
            table_top = table_options.get('top', current_top)
            table_left = table_options.get('left', Inches(1))
            table_width = table_options.get('width', None)
            table_height = table_options.get('height', None)
            column_widths = table_options.get('column_widths', None)
            first_row_is_header = table_options.get('first_row_is_header', True)
            
            add_table(slide,
                     slide_content['table_data'],
                     left=table_left,
                     top=table_top,
                     width=table_width,
                     height=table_height,
                     column_widths=column_widths,
                     first_row_is_header=first_row_is_header)
    
    # Add header/footer
    if 'header' in options or 'footer' in options:
        add_header_footer(prs, 
                         header_text=options.get('header'),
                         footer_text=options.get('footer'),
                         slide_number=options.get('slide_number', True),
                         date=options.get('date', True))
    
    # Save if filename is provided
    if 'filename' in options:
        save_presentation(prs, options['filename'])
    
    return prs